"""
Embeds all page screenshots into the PowerPoint presentation.
Adds one slide per screenshot after the existing content slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, glob

INDIGO     = RGBColor(0x4F, 0x46, 0xE5)
TEAL       = RGBColor(0x0D, 0x94, 0x88)
VIOLET     = RGBColor(0x7C, 0x3A, 0xED)
SUCCESS    = RGBColor(0x05, 0x96, 0x69)
WARNING    = RGBColor(0xD9, 0x77, 0x06)
DARK       = RGBColor(0x1F, 0x29, 0x37)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)

BASE = os.path.dirname(__file__)
PPTX_IN  = os.path.join(BASE, "Project WorkBridge PH User Interface.pptx")
PPTX_OUT = os.path.join(BASE, "Project WorkBridge PH User Interface (with Screenshots).pptx")
SS_DIR   = os.path.join(BASE, "screenshots")

prs = Presentation(PPTX_IN)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width=Pt(0)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = border_width
    else:
        s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h,
             font_size=Pt(12), font_color=DARK, bold=False,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = font_size; r.font.color.rgb = font_color
    r.font.bold = bold; r.font.name = font_name
    return tb

# ── Section separator colors ──────────────────────────────────────────────────
def section_color(label):
    l = label.lower()
    if 'admin'    in l: return VIOLET
    if 'employer' in l: return WARNING
    if 'seeker'   in l: return SUCCESS
    if 'login' in l or 'register' in l or 'forgot' in l: return TEAL
    return INDIGO

# ── Section divider slides ────────────────────────────────────────────────────
SECTIONS = [
    (None,          "01_",  "PUBLIC PAGES",           "Homepage, Jobs, Companies, Safety, About, Help"),
    (None,          "09_",  "AUTHENTICATION PAGES",   "Login, Register, Forgot Password"),
    (None,          "12_",  "JOB SEEKER DASHBOARD",   "Dashboard, Jobs, Applications, Profile, Resume, Notifications, Messages, Settings"),
    (None,          "21_",  "EMPLOYER DASHBOARD",     "Dashboard, Post Job, Manage Jobs, Applicants, Profile, Verification, Messages, Settings"),
    (None,          "29_",  "ADMIN PANEL",            "Dashboard, Verification, Moderation, Reports, Users, Analytics, Categories, Audit Logs, Settings"),
]

section_colors = [INDIGO, TEAL, SUCCESS, WARNING, VIOLET]

screenshots = sorted(glob.glob(os.path.join(SS_DIR, "*.png")))

# Map filename prefix to section
def get_section_color(fname):
    base = os.path.basename(fname)
    num = int(base.split('_')[0])
    if num <= 8:  return INDIGO
    if num <= 11: return TEAL
    if num <= 20: return SUCCESS
    if num <= 28: return WARNING
    return VIOLET

def get_label(fname):
    base = os.path.basename(fname).replace('.png','')
    parts = base.split('_', 1)
    if len(parts) > 1:
        return parts[1].replace('_', ' ').replace('-', ' ').title()
    return base

# ── Add section divider slide ─────────────────────────────────────────────────
def add_section_slide(title, subtitle, color):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, DARK)
    add_rect(slide, 0, 0, Inches(0.12), H, color)
    add_rect(slide, 0, H//2 - Inches(1.2), W, Inches(2.4), color)
    add_text(slide, title,
             Inches(0.5), H//2 - Inches(0.95), W - Inches(1), Inches(0.8),
             font_size=Pt(32), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle,
             Inches(0.5), H//2 + Inches(0.0), W - Inches(1), Inches(0.5),
             font_size=Pt(12), font_color=RGBColor(0xE0,0xE7,0xFF),
             align=PP_ALIGN.CENTER)

# ── Add screenshot slide ──────────────────────────────────────────────────────
def add_screenshot_slide(img_path, label, color):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
    # Top label bar
    add_rect(slide, 0, 0, W, Inches(0.45), color)
    add_text(slide, label,
             Inches(0.2), Inches(0.05), W - Inches(0.4), Inches(0.35),
             font_size=Pt(13), font_color=WHITE, bold=True)
    # Screenshot image
    img_w = W - Inches(0.4)
    img_h = H - Inches(0.65)
    try:
        slide.shapes.add_picture(img_path, Inches(0.2), Inches(0.5), img_w, img_h)
    except Exception as e:
        add_text(slide, f"[Image: {os.path.basename(img_path)}]",
                 Inches(0.5), Inches(2), W - Inches(1), Inches(1),
                 font_size=Pt(14), font_color=GRAY, align=PP_ALIGN.CENTER)
    # Bottom filename label
    add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), DARK)
    add_text(slide, os.path.basename(img_path),
             Inches(0.2), H - Inches(0.26), W - Inches(0.4), Inches(0.22),
             font_size=Pt(8), font_color=GRAY)

# ── Build slides ──────────────────────────────────────────────────────────────
section_triggers = {
    '01': ("PUBLIC PAGES",           "Homepage · Browse Jobs · Job Detail · Companies · Safety · About · Help"),
    '09': ("AUTHENTICATION PAGES",   "Login · Register · Forgot Password"),
    '12': ("JOB SEEKER DASHBOARD",   "Dashboard · Jobs · Saved Jobs · Applications · Profile · Resume · Notifications · Messages · Settings"),
    '21': ("EMPLOYER DASHBOARD",     "Dashboard · Post Job · Manage Jobs · Applicants · Company Profile · Verification · Messages · Settings"),
    '29': ("ADMIN PANEL",            "Dashboard · Verification · Job Moderation · Reports · Users · Analytics · Categories · Audit Logs · Settings"),
}

for img_path in screenshots:
    fname = os.path.basename(img_path)
    prefix = fname.split('_')[0]
    color  = get_section_color(img_path)
    label  = get_label(img_path)

    # Insert section divider before first image of each section
    if prefix in section_triggers:
        title, subtitle = section_triggers[prefix]
        add_section_slide(title, subtitle, color)

    add_screenshot_slide(img_path, label, color)

prs.save(PPTX_OUT)
print(f"Saved: {PPTX_OUT}")
print(f"Total slides: {len(prs.slides)}")

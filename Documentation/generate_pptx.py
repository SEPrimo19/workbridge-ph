"""
WorkBridge PH — MCO Instruction #2 PowerPoint Generator
IT305 – Web Systems and Technologies
Group 4 | Northwest Samar State University
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Brand Colors ────────────────────────────────────────────────────────────
INDIGO       = RGBColor(0x4F, 0x46, 0xE5)   # Primary accent
INDIGO_DARK  = RGBColor(0x43, 0x38, 0xCA)   # Hover
INDIGO_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)   # Muted
TEAL         = RGBColor(0x0D, 0x94, 0x88)   # Secondary
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x1F, 0x29, 0x37)
GRAY         = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY   = RGBColor(0xF3, 0xF4, 0xF6)
SUCCESS      = RGBColor(0x05, 0x96, 0x69)
WARNING      = RGBColor(0xD9, 0x77, 0x06)

# ─── Slide Dimensions (Widescreen 16:9) ──────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

BLANK = prs.slide_layouts[6]  # Completely blank layout


# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=Pt(12), font_color=DARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox

def add_colored_text_box(slide, text, x, y, w, h,
                          bg_color=INDIGO, text_color=WHITE,
                          font_size=Pt(12), bold=False, align=PP_ALIGN.CENTER,
                          border_color=None):
    rect = add_rect(slide, x, y, w, h, bg_color, border_color)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.enum.text import PP_ALIGN as _PP_ALIGN
    tf.paragraphs[0].alignment = align
    # vertical center
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = text_color
    run.font.bold = bold
    run.font.name = "Calibri"
    return rect, txBox

def add_divider(slide, y, color=INDIGO, thickness=Pt(2)):
    line = slide.shapes.add_shape(1, Inches(0.5), y, W - Inches(1), thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def add_bullet_slide_content(slide, items, x, y, w, h, font_size=Pt(11), color=DARK):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = "Calibri"

def slide_header(slide, title, subtitle=None):
    # Top accent bar
    add_rect(slide, 0, 0, W, Inches(0.08), INDIGO)
    # Title
    add_text(slide, title,
             Inches(0.6), Inches(0.2), W - Inches(1.2), Inches(0.55),
             font_size=Pt(22), font_color=DARK, bold=True)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(0.72), W - Inches(1.2), Inches(0.35),
                 font_size=Pt(11), font_color=GRAY)
    # Divider
    add_divider(slide, Inches(1.05))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Full background
add_rect(slide, 0, 0, W, H, DARK)

# Left accent bar
add_rect(slide, 0, 0, Inches(0.12), H, INDIGO)

# Top-left university label
add_text(slide, "NORTHWEST SAMAR STATE UNIVERSITY",
         Inches(0.35), Inches(0.3), Inches(8), Inches(0.35),
         font_size=Pt(9), font_color=GRAY, bold=True)
add_text(slide, "IT305 – Web Systems and Technologies  |  Major Course Output #2",
         Inches(0.35), Inches(0.6), Inches(9), Inches(0.3),
         font_size=Pt(9), font_color=GRAY)

# Main title
add_text(slide, "WorkBridge PH",
         Inches(0.35), Inches(1.5), Inches(10), Inches(1.0),
         font_size=Pt(48), font_color=WHITE, bold=True)

# Indigo underline
add_rect(slide, Inches(0.35), Inches(2.45), Inches(4.5), Inches(0.06), INDIGO)

# Subtitle
add_text(slide, "Web-Based Unified Employment Platform for Decent Work Access",
         Inches(0.35), Inches(2.6), Inches(10), Inches(0.55),
         font_size=Pt(16), font_color=RGBColor(0xD1, 0xD5, 0xDB))

# SDG chips
sdg_labels = ["SDG 8: Decent Work", "SDG 1: No Poverty", "SDG 10: Reduced Inequalities"]
sdg_colors = [INDIGO, TEAL, RGBColor(0x7C, 0x3A, 0xED)]
for i, (lbl, col) in enumerate(zip(sdg_labels, sdg_colors)):
    add_colored_text_box(slide, lbl,
                          Inches(0.35 + i * 3.3), Inches(3.3),
                          Inches(3.0), Inches(0.38),
                          bg_color=col, font_size=Pt(10), bold=True)

# Presenter info
add_text(slide, "Group 4  |  Phase 1 – User Interface",
         Inches(0.35), Inches(4.3), Inches(8), Inches(0.35),
         font_size=Pt(11), font_color=GRAY)
add_text(slide, "April 2026",
         Inches(0.35), Inches(4.6), Inches(4), Inches(0.3),
         font_size=Pt(10), font_color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Project Overview", "What WorkBridge PH is and why it exists")

# Left column — description
add_text(slide, "What is WorkBridge PH?",
         Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4),
         font_size=Pt(13), font_color=INDIGO, bold=True)
add_text(slide,
    "WorkBridge PH is a globally accessible, web-based employment platform that connects "
    "job seekers and employers across all sectors — white-collar, blue-collar, and "
    "household/service workers.\n\n"
    "The platform addresses employment inequality problems that exist worldwide, "
    "using the Philippines as its founding case study.",
    Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.0),
    font_size=Pt(11), font_color=DARK)

# Right column — SDG cards
sdg_data = [
    ("SDG 8", "Decent Work & Economic Growth", INDIGO),
    ("SDG 1", "No Poverty", TEAL),
    ("SDG 10", "Reduced Inequalities", RGBColor(0x7C, 0x3A, 0xED)),
]
for i, (num, label, color) in enumerate(sdg_data):
    bx = Inches(6.8)
    by = Inches(1.3) + Inches(1.1) * i
    add_rect(slide, bx, by, Inches(6.0), Inches(0.9), color)
    add_text(slide, num, bx + Inches(0.15), by + Inches(0.05),
             Inches(1.0), Inches(0.4),
             font_size=Pt(13), font_color=WHITE, bold=True)
    add_text(slide, label, bx + Inches(0.15), by + Inches(0.45),
             Inches(5.5), Inches(0.35),
             font_size=Pt(10), font_color=RGBColor(0xE0, 0xE7, 0xFF))

# Problems section
add_text(slide, "Problems Being Solved",
         Inches(0.5), Inches(3.85), Inches(12), Inches(0.35),
         font_size=Pt(13), font_color=INDIGO, bold=True)
problems = [
    "• Unequal access to digital recruitment platforms for blue-collar and household workers worldwide",
    "• Over-concentration of job portals on white-collar and professional careers only",
    "• Heavy reliance on informal, unsafe recruitment channels and exposure to job scams",
    "• Employers struggling to find reliable workers efficiently across all sectors",
]
add_bullet_slide_content(slide, problems, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5),
                          font_size=Pt(10.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Templates Used
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Templates Used", "Phase 1 — Web & Dashboard Templates")

# Web template card
add_rect(slide, Inches(0.5), Inches(1.25), Inches(5.9), Inches(5.5), LIGHT_GRAY,
         border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(1))
add_colored_text_box(slide, "WEB TEMPLATE (Front-End)",
                      Inches(0.5), Inches(1.25), Inches(5.9), Inches(0.5),
                      bg_color=INDIGO, font_size=Pt(12), bold=True)
items_web = [
    "Framework:  Bootstrap 5.3.2",
    "Icons:  Bootstrap Icons 1.11.1",
    "Styling:  Custom CSS Design Tokens",
    "Scripting:  Vanilla JavaScript (ES Modules)",
    "",
    "Pages Covered:",
    "  • Homepage (index.html)",
    "  • Browse Jobs (jobs.html)",
    "  • Job Detail (job-detail.html)",
    "  • Company Directory (companies.html)",
    "  • Safety & Scam Awareness (safety.html)",
    "  • About, Help, Terms, Privacy",
    "  • Auth: Login, Register, Forgot, Reset",
]
add_bullet_slide_content(slide, items_web,
                          Inches(0.7), Inches(1.85), Inches(5.5), Inches(4.5),
                          font_size=Pt(10.5))

# Dashboard template card
add_rect(slide, Inches(6.9), Inches(1.25), Inches(5.9), Inches(5.5), LIGHT_GRAY,
         border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(1))
add_colored_text_box(slide, "DASHBOARD TEMPLATE (Backend UI)",
                      Inches(6.9), Inches(1.25), Inches(5.9), Inches(0.5),
                      bg_color=TEAL, font_size=Pt(12), bold=True)
items_dash = [
    "Layout:  Sidebar + Topbar + Content Area",
    "Framework:  Bootstrap 5.3.2",
    "Styling:  Custom role-specific CSS",
    "Dark Mode:  CSS data-theme toggle",
    "",
    "Dashboards Covered:",
    "  • Job Seeker Dashboard (10 pages)",
    "  • Employer Dashboard (10 pages)",
    "  • Admin Panel (13 pages)",
    "",
    "Common Components:",
    "  • Stat cards, data tables, modals",
    "  • Toast notifications, dropdowns",
    "  • Form pages with JS validation",
]
add_bullet_slide_content(slide, items_dash,
                          Inches(7.1), Inches(1.85), Inches(5.5), Inches(4.5),
                          font_size=Pt(10.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — UI Component: Homepage
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Required UI Component — Homepage", "index.html")

components = [
    ("Website Header",     "Public navigation bar with WorkBridge PH brand logo and navigation links",              "✔"),
    ("Navigation Menu",    "Links: Jobs, Companies, About, Help — plus Login button",                               "✔"),
    ("Banner Section",     "Hero section with headline, subtitle, and job search form (keyword + location)",        "✔"),
    ("System Introduction","'How It Works' (3 steps) + 'Browse by Category' + 'Featured Jobs' grid sections",     "✔"),
    ("Footer",             "Copyright notice © 2026 WorkBridge PH with Safety, Terms, Privacy links",              "✔"),
]

# Table header
cols = [Inches(2.5), Inches(7.5), Inches(2.0)]
xs   = [Inches(0.5), Inches(3.0), Inches(10.6)]
add_colored_text_box(slide, "Component", xs[0], Inches(1.25), cols[0], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Description", xs[1], Inches(1.25), cols[1], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Status", xs[2], Inches(1.25), cols[2], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)

for i, (name, desc, status) in enumerate(components):
    row_y = Inches(1.7) + Inches(0.72) * i
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    for j, (tx, col) in enumerate(zip(xs, cols)):
        add_rect(slide, tx, row_y, col, Inches(0.65), bg,
                 border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(0.5))
    vals = [name, desc, status]
    colors = [DARK, GRAY, SUCCESS]
    bolds  = [True, False, True]
    for j, (tx, col, val, col_c, bld) in enumerate(zip(xs, cols, vals, colors, bolds)):
        add_text(slide, val, tx + Inches(0.1), row_y + Inches(0.1),
                 col - Inches(0.2), Inches(0.5),
                 font_size=Pt(10), font_color=col_c, bold=bld)

add_text(slide, "✔ All 5 required homepage components are implemented.",
         Inches(0.5), Inches(5.6), Inches(12), Inches(0.4),
         font_size=Pt(11), font_color=SUCCESS, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — UI Component: Login Page
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Required UI Component — Login Page", "auth/login.html")

components = [
    ("Username Field",          "Email input field (type='email') with placeholder",                                    "✔"),
    ("Password Field",          "Password input field (type='password') with show/hide toggle",                        "✔"),
    ("Login Button",            "Submit button — triggers mock authentication on click",                                "✔"),
    ("JavaScript Validation",   "auth.mock.js validates credentials; toastError() shown on failure; role-based redirect on success", "✔"),
]

xs   = [Inches(0.5), Inches(3.0), Inches(10.6)]
cols = [Inches(2.5), Inches(7.5), Inches(2.0)]

add_colored_text_box(slide, "Component", xs[0], Inches(1.25), cols[0], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Description", xs[1], Inches(1.25), cols[1], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Status", xs[2], Inches(1.25), cols[2], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)

for i, (name, desc, status) in enumerate(components):
    row_y = Inches(1.7) + Inches(0.72) * i
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    for j, (tx, col) in enumerate(zip(xs, cols)):
        add_rect(slide, tx, row_y, col, Inches(0.65), bg,
                 border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(0.5))
    vals   = [name, desc, status]
    colors = [DARK, GRAY, SUCCESS]
    bolds  = [True, False, True]
    for j, (tx, col, val, col_c, bld) in enumerate(zip(xs, cols, vals, colors, bolds)):
        add_text(slide, val, tx + Inches(0.1), row_y + Inches(0.1),
                 col - Inches(0.2), Inches(0.5),
                 font_size=Pt(10), font_color=col_c, bold=bld)

# Demo accounts note
add_rect(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.5),
         INDIGO_LIGHT, border_color=INDIGO, border_width=Pt(1))
add_text(slide, "Demo Accounts (for testing)",
         Inches(0.7), Inches(4.95), Inches(6), Inches(0.3),
         font_size=Pt(11), font_color=INDIGO, bold=True)
accounts = [
    "Job Seeker:  seeker@demo.com  /  demo123    →  Redirects to Seeker Dashboard",
    "Employer:    employer@demo.com  /  demo123  →  Redirects to Employer Dashboard",
    "Admin:       admin@demo.com  /  demo123     →  Redirects to Admin Dashboard",
]
add_bullet_slide_content(slide, accounts,
                          Inches(0.7), Inches(5.25), Inches(12), Inches(1.0),
                          font_size=Pt(10), color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — UI Component: Dashboard Page
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Required UI Component — Dashboard Page",
             "seeker/dashboard.html  |  employer/dashboard.html  |  admin/dashboard.html")

components = [
    ("Navigation Sidebar", "Role-specific sidebar with icon + label nav links, rendered by layout.js for all 3 roles",  "✔"),
    ("Summary Info Cards", "Stat cards showing counts (saved jobs, applications, profile strength, pending queues)",     "✔"),
    ("Content Area",       ".app-content div — main working area with charts, tables, and role-specific data widgets",  "✔"),
]

xs   = [Inches(0.5), Inches(3.0), Inches(10.6)]
cols = [Inches(2.5), Inches(7.5), Inches(2.0)]

add_colored_text_box(slide, "Component", xs[0], Inches(1.25), cols[0], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Description", xs[1], Inches(1.25), cols[1], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Status", xs[2], Inches(1.25), cols[2], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)

for i, (name, desc, status) in enumerate(components):
    row_y = Inches(1.7) + Inches(0.72) * i
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    for j, (tx, col) in enumerate(zip(xs, cols)):
        add_rect(slide, tx, row_y, col, Inches(0.65), bg,
                 border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(0.5))
    vals   = [name, desc, status]
    colors = [DARK, GRAY, SUCCESS]
    bolds  = [True, False, True]
    for j, (tx, col, val, col_c, bld) in enumerate(zip(xs, cols, vals, colors, bolds)):
        add_text(slide, val, tx + Inches(0.1), row_y + Inches(0.1),
                 col - Inches(0.2), Inches(0.5),
                 font_size=Pt(10), font_color=col_c, bold=bld)

# Three dashboard cards
dash_data = [
    ("Job Seeker Dashboard", ["Saved jobs count", "Applications count", "Profile strength %", "Recommended jobs grid"], INDIGO),
    ("Employer Dashboard",   ["Active/Pending jobs", "Total applicants", "Applicant chart", "Quick post job"], TEAL),
    ("Admin Dashboard",      ["Pending verifications", "Pending job reviews", "Open reports", "Analytics overview"], RGBColor(0x7C, 0x3A, 0xED)),
]
for i, (title, items, color) in enumerate(dash_data):
    bx = Inches(0.5 + i * 4.25)
    by = Inches(4.05)
    add_colored_text_box(slide, title, bx, by, Inches(4.0), Inches(0.4),
                          bg_color=color, font_size=Pt(10), bold=True)
    add_rect(slide, bx, by + Inches(0.4), Inches(4.0), Inches(2.5), LIGHT_GRAY,
             border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(0.5))
    add_bullet_slide_content(slide, ["• " + x for x in items],
                              bx + Inches(0.1), by + Inches(0.5),
                              Inches(3.8), Inches(2.2), font_size=Pt(10))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — UI Component: Data Entry Form Page
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Required UI Component — Data Entry Form Page",
             "employer/post-job.html  |  auth/register.html  |  seeker/profile.html")

components = [
    ("Input Fields",         "Text, select, textarea, number inputs — title, type, setup, location, salary, description", "✔"),
    ("Labels",               "All fields are labelled with <label> tags and placeholder text",                              "✔"),
    ("Submit Button",        "Primary submit button triggers form processing and localStorage/API write",                   "✔"),
    ("JavaScript Validation","HTML5 required attributes + JS checks — error toasts shown for invalid/missing fields",      "✔"),
]

xs   = [Inches(0.5), Inches(3.0), Inches(10.6)]
cols = [Inches(2.5), Inches(7.5), Inches(2.0)]

add_colored_text_box(slide, "Component", xs[0], Inches(1.25), cols[0], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Description", xs[1], Inches(1.25), cols[1], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Status", xs[2], Inches(1.25), cols[2], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)

for i, (name, desc, status) in enumerate(components):
    row_y = Inches(1.7) + Inches(0.72) * i
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    for j, (tx, col) in enumerate(zip(xs, cols)):
        add_rect(slide, tx, row_y, col, Inches(0.65), bg,
                 border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(0.5))
    vals   = [name, desc, status]
    colors = [DARK, GRAY, SUCCESS]
    bolds  = [True, False, True]
    for j, (tx, col, val, col_c, bld) in enumerate(zip(xs, cols, vals, colors, bolds)):
        add_text(slide, val, tx + Inches(0.1), row_y + Inches(0.1),
                 col - Inches(0.2), Inches(0.5),
                 font_size=Pt(10), font_color=col_c, bold=bld)

# Form pages list
add_rect(slide, Inches(0.5), Inches(4.75), Inches(12.3), Inches(2.0),
         INDIGO_LIGHT, border_color=INDIGO, border_width=Pt(1))
add_text(slide, "Data Entry Form Pages in This System",
         Inches(0.7), Inches(4.85), Inches(8), Inches(0.3),
         font_size=Pt(11), font_color=INDIGO, bold=True)
form_pages = [
    "employer/post-job.html — Post a Job form (title, type, setup, location, salary, description, requirements, benefits)",
    "auth/register.html — User Registration form with role selection (Seeker / Employer)",
    "auth/login.html — Login form with email + password validation",
    "seeker/profile.html — Editable profile form (name, location, skills) via modal",
]
add_bullet_slide_content(slide, ["• " + p for p in form_pages],
                          Inches(0.7), Inches(5.15), Inches(12.1), Inches(1.5),
                          font_size=Pt(10))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — UI Component: Data Display Page
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Required UI Component — Data Display Page",
             "employer/manage-jobs.html  |  admin/users.html  |  admin/categories.html")

components = [
    ("Table Layout",   "Responsive HTML tables with thead/tbody, status badges, and pagination chips",          "✔"),
    ("Edit Button",    "Links to edit page (job-edit.html, user-detail.html) for modifying records",            "✔"),
    ("View Button",    "Links to detail page (job-detail.html, user-detail.html) for viewing full record",      "✔"),
    ("Delete Button",  "Confirmation prompt → removes record from data store; success toast shown after delete","✔"),
]

xs   = [Inches(0.5), Inches(3.0), Inches(10.6)]
cols = [Inches(2.5), Inches(7.5), Inches(2.0)]

add_colored_text_box(slide, "Component", xs[0], Inches(1.25), cols[0], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Description", xs[1], Inches(1.25), cols[1], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_colored_text_box(slide, "Status", xs[2], Inches(1.25), cols[2], Inches(0.4),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)

for i, (name, desc, status) in enumerate(components):
    row_y = Inches(1.7) + Inches(0.72) * i
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    for j, (tx, col) in enumerate(zip(xs, cols)):
        add_rect(slide, tx, row_y, col, Inches(0.65), bg,
                 border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(0.5))
    vals   = [name, desc, status]
    colors = [DARK, GRAY, SUCCESS]
    bolds  = [True, False, True]
    for j, (tx, col, val, col_c, bld) in enumerate(zip(xs, cols, vals, colors, bolds)):
        add_text(slide, val, tx + Inches(0.1), row_y + Inches(0.1),
                 col - Inches(0.2), Inches(0.5),
                 font_size=Pt(10), font_color=col_c, bold=bld)

# Display pages list
add_rect(slide, Inches(0.5), Inches(4.75), Inches(12.3), Inches(2.0),
         INDIGO_LIGHT, border_color=INDIGO, border_width=Pt(1))
add_text(slide, "Data Display Pages in This System",
         Inches(0.7), Inches(4.85), Inches(8), Inches(0.3),
         font_size=Pt(11), font_color=INDIGO, bold=True)
display_pages = [
    "employer/manage-jobs.html — Job listings table with Edit, View, Delete buttons",
    "admin/users.html — User management table with Edit, View, Delete buttons",
    "admin/categories.html — Job categories table with Delete button",
    "seeker/applications.html — Application status cards with filter chips",
    "admin/employer-verification.html — Employer verification table with status badges",
]
add_bullet_slide_content(slide, ["• " + p for p in display_pages],
                          Inches(0.7), Inches(5.15), Inches(12.1), Inches(1.5),
                          font_size=Pt(10))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Full UI Summary
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Complete UI Summary", "All pages implemented in Phase 1")

sections = [
    ("Public Pages", [
        "index.html — Homepage",
        "jobs.html — Browse Jobs",
        "job-detail.html — Job Detail",
        "companies.html — Company Directory",
        "company-detail.html — Company Detail",
        "safety.html — Scam Awareness",
        "about.html / help.html",
        "terms.html / privacy.html",
    ], INDIGO),
    ("Auth Pages", [
        "auth/login.html",
        "auth/register.html",
        "auth/forgot.html",
        "auth/reset.html",
    ], TEAL),
    ("Seeker Dashboard", [
        "dashboard.html",
        "jobs.html",
        "saved-jobs.html",
        "applications.html",
        "profile.html",
        "resume.html",
        "notifications.html",
        "messages.html",
        "settings.html",
    ], RGBColor(0x05, 0x96, 0x69)),
    ("Employer Dashboard", [
        "dashboard.html",
        "post-job.html",
        "manage-jobs.html",
        "applicants.html",
        "company-profile.html",
        "verification.html",
        "messages.html",
        "settings.html",
    ], RGBColor(0xD9, 0x77, 0x06)),
    ("Admin Panel", [
        "dashboard.html",
        "employer-verification.html",
        "job-moderation.html",
        "reports.html",
        "users.html",
        "analytics.html",
        "categories.html",
        "audit-logs.html",
    ], RGBColor(0x7C, 0x3A, 0xED)),
]

col_w = Inches(2.45)
for i, (title, pages, color) in enumerate(sections):
    bx = Inches(0.3) + col_w * i
    by = Inches(1.2)
    add_colored_text_box(slide, title, bx, by, col_w - Inches(0.1), Inches(0.4),
                          bg_color=color, font_size=Pt(9.5), bold=True)
    add_rect(slide, bx, by + Inches(0.4), col_w - Inches(0.1), Inches(5.2),
             LIGHT_GRAY, border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(0.5))
    add_bullet_slide_content(slide, ["• " + p for p in pages],
                              bx + Inches(0.1), by + Inches(0.5),
                              col_w - Inches(0.2), Inches(5.0),
                              font_size=Pt(9))

# Total count
add_text(slide, "Total: 33 pages  |  HTML + CSS + JavaScript  |  Bootstrap 5.3.2  |  Dark Mode Support",
         Inches(0.3), Inches(7.05), Inches(12), Inches(0.35),
         font_size=Pt(10), font_color=GRAY, bold=False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Technology Stack
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, WHITE)
slide_header(slide, "Technology Stack", "Phase 1 (Current) and Phase 2 (Planned)")

# Phase 1
add_colored_text_box(slide, "PHASE 1 — Frontend Prototype (Current)",
                      Inches(0.5), Inches(1.25), Inches(5.9), Inches(0.45),
                      bg_color=INDIGO, font_size=Pt(11), bold=True)
add_rect(slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(4.8), LIGHT_GRAY,
         border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(1))
phase1 = [
    "HTML5 — Page structure",
    "CSS3 Custom Properties — Design tokens & theming",
    "Vanilla JavaScript (ES Modules) — App logic",
    "Bootstrap 5.3.2 — Responsive UI framework",
    "Bootstrap Icons 1.11.1 — Iconography",
    "localStorage — Mock data backend",
    "Custom client-side Router — Role-based guards",
    "Playwright — End-to-end testing",
]
add_bullet_slide_content(slide, ["• " + x for x in phase1],
                          Inches(0.7), Inches(1.85), Inches(5.5), Inches(4.5),
                          font_size=Pt(10.5))

# Phase 2
add_colored_text_box(slide, "PHASE 2 — Full-Stack Target (Next)",
                      Inches(6.9), Inches(1.25), Inches(5.9), Inches(0.45),
                      bg_color=TEAL, font_size=Pt(11), bold=True)
add_rect(slide, Inches(6.9), Inches(1.7), Inches(5.9), Inches(4.8), LIGHT_GRAY,
         border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(1))
phase2 = [
    "PHP 8.2+ — Server-side language",
    "Laravel 11.x — MVC application framework",
    "Laravel Sanctum — API token authentication",
    "Eloquent ORM — Database models & relations",
    "Laravel Form Requests — Server-side validation",
    "Laravel Gates & Policies — Authorization",
    "Laravel Storage — File uploads (resumes, docs)",
    "MySQL 8.0+ — Relational database",
    "PHPUnit — Backend unit & feature testing",
    "XAMPP — Local development environment",
]
add_bullet_slide_content(slide, ["• " + x for x in phase2],
                          Inches(7.1), Inches(1.85), Inches(5.5), Inches(4.5),
                          font_size=Pt(10.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Closing Slide
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, DARK)
add_rect(slide, 0, 0, Inches(0.12), H, INDIGO)

add_text(slide, "WorkBridge PH",
         Inches(0.5), Inches(2.2), Inches(12), Inches(1.0),
         font_size=Pt(42), font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.5), Inches(3.1), Inches(4.33), Inches(0.06), INDIGO)

add_text(slide, "Phase 1 — User Interface Complete",
         Inches(0.5), Inches(3.25), Inches(12), Inches(0.5),
         font_size=Pt(16), font_color=RGBColor(0xD1, 0xD5, 0xDB), align=PP_ALIGN.CENTER)

add_text(slide, "IT305 – Web Systems and Technologies  |  Group 4",
         Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
         font_size=Pt(11), font_color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "Northwest Samar State University  |  April 2026",
         Inches(0.5), Inches(4.4), Inches(12), Inches(0.4),
         font_size=Pt(10), font_color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "Thank You",
         Inches(0.5), Inches(5.4), Inches(12), Inches(0.6),
         font_size=Pt(28), font_color=INDIGO, bold=True, align=PP_ALIGN.CENTER)


# ─── Save ─────────────────────────────────────────────────────────────────────
import os
out_path = os.path.join(os.path.dirname(__file__), "Project WorkBridge PH User Interface.pptx")
prs.save(out_path)
print("Saved:", out_path)
